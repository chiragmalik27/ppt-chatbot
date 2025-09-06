import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import plotly.express as px
import plotly.graph_objects as go
import pandas as pd
import matplotlib.pyplot as plt
from io import BytesIO
import os
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configure Gemini AI
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

class PowerPointChatbot:
    def __init__(self):
        self.model = genai.GenerativeModel('gemini-2.5-flash')
        self.current_ppt = None
        
    def create_presentation(self, title, content_structure):
        """Create a new PowerPoint presentation with professional design"""
        prs = Presentation()
        
        # Create title slide with enhanced styling
        self.create_title_slide(prs, title)
        
        # Add content slides based on structure
        for i, slide_data in enumerate(content_structure):
            self.add_content_slide(prs, slide_data, i)
        
        # Add conclusion slide
        self.add_conclusion_slide(prs, title)
            
        self.current_ppt = prs
        return prs
    
    def create_title_slide(self, prs, title):
        """Create an enhanced title slide with professional design"""
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Set slide background to professional blue gradient
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = RGBColor(41, 98, 255)  # Light blue
        fill.gradient_stops[1].color.rgb = RGBColor(0, 56, 168)   # Dark blue
        
        # Set title with enhanced styling
        title_shape = slide.shapes.title
        title_shape.text = title.upper()  # Make title uppercase for impact
        
        # Style the title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(48)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Position title higher on slide
        title_shape.top = Inches(2.5)
        title_shape.left = Inches(1)
        title_shape.width = Inches(8)
        title_shape.height = Inches(2)
        
        # Add subtitle with professional styling
        if len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = "AI-Powered Presentation"
            subtitle_paragraph = subtitle_shape.text_frame.paragraphs[0]
            subtitle_paragraph.font.size = Pt(24)
            subtitle_paragraph.font.color.rgb = RGBColor(200, 220, 255)  # Light blue text
            subtitle_paragraph.alignment = PP_ALIGN.CENTER
            
            # Position subtitle
            subtitle_shape.top = Inches(5)
            subtitle_shape.left = Inches(1)
            subtitle_shape.width = Inches(8)
        
        # Add decorative elements - blue accent bars
        self.add_decorative_elements_title(slide)
    
    def add_decorative_elements_title(self, slide):
        """Add decorative elements to title slide"""
        # Add geometric shapes similar to the second image
        # Left side decorative elements
        for i in range(3):
            left = Inches(0.5 + i * 0.3)
            top = Inches(1 + i * 0.8)
            width = Inches(0.2)
            height = Inches(0.2)
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                left, top, width, height
            )
            
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(100, 200, 255)  # Light blue circles
            line = shape.line
            line.fill.background()
        
        # Right side decorative elements
        for i in range(3):
            left = Inches(8.5 + i * 0.3)
            top = Inches(1.5 + i * 0.6)
            width = Inches(0.15)
            height = Inches(0.15)
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                left, top, width, height
            )
            
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(150, 220, 255)  # Lighter blue circles
            line = shape.line
            line.fill.background()

    def add_content_slide(self, prs, slide_data, slide_index=0):
        """Add a professionally styled content slide matching the blue theme"""
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide background with blue theme
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = RGBColor(240, 248, 255)  # Very light blue
        fill.gradient_stops[1].color.rgb = RGBColor(220, 240, 255)  # Light blue
        
        # Set title with enhanced styling to match theme
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get('title', 'Slide Title').upper()
        
        # Style the title with blue theme
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(40)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 56, 168)  # Dark blue
        title_paragraph.alignment = PP_ALIGN.LEFT
        
        # Position title (adjusted to ensure good spacing)
        title_shape.top = Inches(0.8)
        title_shape.left = Inches(0.6)  # Moved slightly right to clear the blue bar
        title_shape.width = Inches(8.8)  # Made wider to use available space
        title_shape.height = Inches(1.2)
        
        # Add content using the content placeholder (this is the key fix)
        content = slide_data.get('content', [])
        if content and len(content) > 0:
            # Use the built-in content placeholder
            content_placeholder = None
            for shape in slide.shapes:
                if shape.is_placeholder:
                    placeholder_format = shape.placeholder_format
                    if placeholder_format.idx == 1:  # Content placeholder
                        content_placeholder = shape
                        break
            
            if content_placeholder:
                # Adjust content placeholder position to avoid blue bar overlap
                content_placeholder.left = Inches(0.6)  # Move right to clear blue bar
                content_placeholder.top = Inches(2.2)
                content_placeholder.width = Inches(8.8)  # Use more available space
                content_placeholder.height = Inches(4.5)
                
                # Clear and set up the text frame
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                # Set text frame properties for better formatting
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                text_frame.word_wrap = True
                
                # Add each content point as a separate paragraph
                for i, point in enumerate(content):
                    # Create paragraph
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    # Clean up the content point
                    clean_point = str(point).strip()
                    if clean_point:
                        # Ensure proper bullet formatting
                        if not clean_point.startswith(('•', '-', '*')):
                            clean_point = f"• {clean_point}"
                        
                        # Set the text
                        p.text = clean_point
                        p.level = 0
                        
                        # Apply formatting
                        p.font.size = Pt(18)
                        p.font.name = 'Calibri'
                        p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray for readability
                        p.space_after = Pt(12)
                        p.space_before = Pt(6)
                        p.alignment = PP_ALIGN.LEFT
                        p.line_spacing = 1.15
            else:
                # Fallback: create a new text box if placeholder not found
                left = Inches(0.6)  # Adjusted to clear the blue bar
                top = Inches(2.2)
                width = Inches(8.8)  # Made wider to use available space
                height = Inches(4.5)
                
                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                text_frame.clear()
                
                # Add content to text box
                for i, point in enumerate(content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    clean_point = str(point).strip()
                    if clean_point:
                        if not clean_point.startswith(('•', '-', '*')):
                            clean_point = f"• {clean_point}"
                        
                        p.text = clean_point
                        p.level = 0
                        p.font.size = Pt(18)
                        p.font.name = 'Calibri'
                        p.font.color.rgb = RGBColor(51, 51, 51)
                        p.space_after = Pt(12)
                        p.alignment = PP_ALIGN.LEFT
        
        # Add decorative blue accent bar (moved further left to avoid text overlap)
        left = Inches(0.2)  # Moved much further left
        top = Inches(0.8)
        width = Inches(0.2)  # Made slightly wider for better visual impact
        height = Inches(6)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        
        # Style the decorative element with blue gradient
        fill = shape.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = RGBColor(41, 98, 255)  # Light blue
        fill.gradient_stops[1].color.rgb = RGBColor(0, 56, 168)   # Dark blue
        
        # Remove shape outline
        line = shape.line
        line.fill.background()
        
        # Add flowing decorative elements (circles and lines)
        self.add_decorative_elements_content(slide)
    
    def add_decorative_elements_content(self, slide):
        """Add flowing decorative elements to content slides"""
        # Add flowing line elements (like in the second image)
        import random
        
        # Right side decorative flowing elements
        for i in range(5):
            left = Inches(8.5 + random.uniform(-0.2, 0.3))
            top = Inches(1 + i * 1.2 + random.uniform(-0.3, 0.3))
            width = Inches(0.1 + random.uniform(0, 0.1))
            height = Inches(0.1 + random.uniform(0, 0.1))
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                left, top, width, height
            )
            
            fill = shape.fill
            fill.solid()
            alpha_values = [100, 150, 200]
            blue_value = random.choice(alpha_values)
            fill.fore_color.rgb = RGBColor(blue_value, blue_value + 50, 255)
            line = shape.line
            line.fill.background()
        
        # Add connecting line elements
        for i in range(3):
            left = Inches(8.2)
            top = Inches(2 + i * 1.5)
            width = Inches(1.5)
            height = Inches(0.05)
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )
            
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(150, 200, 255)  # Light blue lines
            line = shape.line
            line.fill.background()

    def add_conclusion_slide(self, prs, title):
        """Add a conclusion slide with blue theme"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide background to professional blue gradient (same as title)
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = RGBColor(41, 98, 255)  # Light blue
        fill.gradient_stops[1].color.rgb = RGBColor(0, 56, 168)   # Dark blue
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "THANK YOU!"
        
        # Style the title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(54)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Position title
        title_shape.top = Inches(2.5)
        title_shape.left = Inches(1)
        title_shape.width = Inches(8)
        title_shape.height = Inches(2)
        
        # Add subtitle
        if len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = f"Questions & Discussion\n\nPresentation: {title}"
            subtitle_paragraph = subtitle_shape.text_frame.paragraphs[0]
            subtitle_paragraph.font.size = Pt(24)
            subtitle_paragraph.font.color.rgb = RGBColor(200, 220, 255)  # Light blue
            subtitle_paragraph.alignment = PP_ALIGN.CENTER
            
            # Position subtitle
            subtitle_shape.top = Inches(5)
            subtitle_shape.left = Inches(1)
            subtitle_shape.width = Inches(8)
        
        # Add decorative elements
        self.add_decorative_elements_title(slide)
    
    def edit_slide_content(self, slide_index, new_content):
        """Edit the content of a specific slide"""
        if not self.current_ppt:
            return False
            
        try:
            slide = self.current_ppt.slides[slide_index]
            
            # Update title if provided
            if 'title' in new_content:
                slide.shapes.title.text = new_content['title']
            
            # Update content if provided
            if 'content' in new_content:
                # Find the content placeholder or text frame
                content_updated = False
                
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        text_frame = shape.text_frame
                        text_frame.clear()
                        
                        # Add new content points
                        for i, point in enumerate(new_content['content']):
                            if point.strip():  # Only add non-empty content
                                if i == 0:
                                    p = text_frame.paragraphs[0]
                                else:
                                    p = text_frame.add_paragraph()
                                
                                # Clean and format the point
                                clean_point = str(point).strip()
                                if not clean_point.startswith(('•', '-', '*')):
                                    clean_point = f"• {clean_point}"
                                
                                p.text = clean_point
                                p.level = 0
                                
                                # Apply formatting
                                p.font.size = Pt(18)
                                p.font.name = 'Calibri'
                                p.font.color.rgb = RGBColor(51, 51, 51)
                        
                        content_updated = True
                        break
                
                # If no text frame found, create a new text box
                if not content_updated:
                    left = Inches(0.6)
                    top = Inches(2.2)
                    width = Inches(8.8)
                    height = Inches(4.5)
                    
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    
                    for i, point in enumerate(new_content['content']):
                        if point.strip():
                            if i == 0:
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            
                            clean_point = str(point).strip()
                            if not clean_point.startswith(('•', '-', '*')):
                                clean_point = f"• {clean_point}"
                            
                            p.text = clean_point
                            p.level = 0
                            p.font.size = Pt(18)
                            p.font.name = 'Calibri'
                            p.font.color.rgb = RGBColor(51, 51, 51)
            
            return True
        except (IndexError, Exception) as e:
            print(f"Error editing slide: {e}")
            return False

    def add_new_slide(self, slide_content):
        """Add a new slide to the existing presentation"""
        if not self.current_ppt:
            return False
        
        try:
            # Use the content slide layout
            slide_layout = self.current_ppt.slide_layouts[1]  # Title and Content layout
            slide = self.current_ppt.slides.add_slide(slide_layout)
            
            # Set slide background with blue theme
            background = slide.background
            fill = background.fill
            fill.gradient()
            fill.gradient_stops[0].color.rgb = RGBColor(240, 248, 255)  # Very light blue
            fill.gradient_stops[1].color.rgb = RGBColor(220, 240, 255)  # Light blue
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = slide_content.get('title', 'New Slide').upper()
            
            # Style the title
            title_paragraph = title_shape.text_frame.paragraphs[0]
            title_paragraph.font.size = Pt(40)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = RGBColor(0, 56, 168)  # Dark blue
            title_paragraph.alignment = PP_ALIGN.LEFT
            
            # Position title
            title_shape.top = Inches(0.8)
            title_shape.left = Inches(0.6)
            title_shape.width = Inches(8.8)
            title_shape.height = Inches(1.2)
            
            # Add content
            content = slide_content.get('content', [])
            if content:
                # Find content placeholder
                content_placeholder = None
                for shape in slide.shapes:
                    if shape.is_placeholder:
                        placeholder_format = shape.placeholder_format
                        if placeholder_format.idx == 1:  # Content placeholder
                            content_placeholder = shape
                            break
                
                if content_placeholder:
                    # Adjust content placeholder position
                    content_placeholder.left = Inches(0.6)
                    content_placeholder.top = Inches(2.2)
                    content_placeholder.width = Inches(8.8)
                    content_placeholder.height = Inches(4.5)
                    
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    # Add content points
                    for i, point in enumerate(content):
                        if point.strip():
                            if i == 0:
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            
                            clean_point = str(point).strip()
                            if not clean_point.startswith(('•', '-', '*')):
                                clean_point = f"• {clean_point}"
                            
                            p.text = clean_point
                            p.level = 0
                            p.font.size = Pt(18)
                            p.font.name = 'Calibri'
                            p.font.color.rgb = RGBColor(51, 51, 51)
                            p.space_after = Pt(12)
                            p.alignment = PP_ALIGN.LEFT
            
            # Add decorative blue accent bar
            left = Inches(0.2)
            top = Inches(0.8)
            width = Inches(0.2)
            height = Inches(5.5)
            
            accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            fill = accent_bar.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(41, 98, 255)  # Blue accent
            accent_bar.line.fill.background()
            
            return True
        except Exception as e:
            print(f"Error adding new slide: {e}")
            return False
    
    def add_chart_slide(self, chart_data, chart_type="bar"):
        """Add a slide with a chart"""
        if not self.current_ppt:
            return False
        
        slide_layout = self.current_ppt.slide_layouts[5]  # Blank layout
        slide = self.current_ppt.slides.add_slide(slide_layout)
        
        # Create chart using matplotlib
        fig, ax = plt.subplots(figsize=(10, 6))
        
        if chart_type == "bar":
            ax.bar(chart_data['labels'], chart_data['values'])
        elif chart_type == "line":
            ax.plot(chart_data['labels'], chart_data['values'])
        elif chart_type == "pie":
            ax.pie(chart_data['values'], labels=chart_data['labels'], autopct='%1.1f%%')
        
        ax.set_title(chart_data.get('title', 'Chart'))
        
        # Save chart as image
        img_buffer = BytesIO()
        plt.savefig(img_buffer, format='png', bbox_inches='tight', dpi=300)
        img_buffer.seek(0)
        plt.close()
        
        # Add image to slide
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(6)
        
        slide.shapes.add_picture(img_buffer, left, top, width, height)
        
        return True
    
    def generate_content_with_ai(self, prompt):
        """Generate content using Gemini AI"""
        try:
            response = self.model.generate_content(prompt)
            return response.text
        except Exception as e:
            return f"Error generating content: {str(e)}"
    
    def extract_topic_from_prompt(self, prompt):
        """Extract the main topic from user's presentation request."""
        # Keywords that typically indicate the topic
        topic_indicators = [
            "presentation on", "presentation about", "about", "on", 
            "slides about", "slides on", "ppt on", "ppt about",
            "create a presentation on", "create a presentation about",
            "make a presentation on", "make a presentation about"
        ]
        
        # Clean the prompt
        clean_prompt = prompt.lower().strip()
        
        # Find the topic after common indicators
        for indicator in topic_indicators:
            if indicator in clean_prompt:
                # Split by the indicator and take the part after it
                parts = clean_prompt.split(indicator, 1)
                if len(parts) > 1:
                    topic = parts[1].strip()
                    # Remove common words at the beginning
                    topic = topic.replace("the ", "").replace("a ", "").replace("an ", "")
                    # Remove trailing requests like "and i want 5 slides"
                    if " and " in topic:
                        topic = topic.split(" and ")[0]
                    if " with " in topic:
                        topic = topic.split(" with ")[0]
                    return topic.strip()
        
        # If no specific indicator found, try to extract from common patterns
        if "create" in clean_prompt or "make" in clean_prompt or "generate" in clean_prompt:
            # Look for words after "presentation" or "ppt"
            words = clean_prompt.split()
            for i, word in enumerate(words):
                if word in ["presentation", "ppt", "powerpoint", "slides"]:
                    if i + 1 < len(words) and words[i + 1] in ["on", "about"]:
                        # Get everything after "on/about"
                        remaining = " ".join(words[i + 2:])
                        if " and " in remaining:
                            remaining = remaining.split(" and ")[0]
                        if " with " in remaining:
                            remaining = remaining.split(" with ")[0]
                        return remaining.strip()
        
        # Fallback: return a cleaned version of the prompt
        fallback = clean_prompt.replace("create", "").replace("make", "").replace("generate", "")
        fallback = fallback.replace("presentation", "").replace("ppt", "").replace("powerpoint", "")
        fallback = fallback.replace("a ", "").replace("the ", "").strip()
        
        return fallback if fallback else "General Topic"

    def extract_slide_count_from_prompt(self, prompt):
        """Extract the number of slides requested from user's prompt."""
        import re
        
        # Look for patterns like "5 slides", "include 5 slides", "with 5 slides", etc.
        patterns = [
            r'(\d+)\s*slides?',
            r'include\s*(\d+)\s*slides?',
            r'with\s*(\d+)\s*slides?',
            r'want\s*(\d+)\s*slides?',
            r'need\s*(\d+)\s*slides?',
            r'make\s*(\d+)\s*slides?'
        ]
        
        for pattern in patterns:
            match = re.search(pattern, prompt.lower())
            if match:
                count = int(match.group(1))
                # Reasonable range check
                if 2 <= count <= 15:
                    return count
        
        # Default to 5 slides if no specific count found
        return 5

    def extract_slide_number_from_prompt(self, prompt):
        """Extract slide number from user's editing request."""
        import re
        
        # Look for patterns like "slide 1", "slide number 2", "first slide", etc.
        patterns = [
            r'slide\s*(\d+)',
            r'slide\s*number\s*(\d+)',
            r'(\d+)(?:st|nd|rd|th)?\s*slide',
        ]
        
        for pattern in patterns:
            match = re.search(pattern, prompt.lower())
            if match:
                return int(match.group(1))
        
        # Look for ordinal words
        ordinals = {
            'first': 1, 'second': 2, 'third': 3, 'fourth': 4, 'fifth': 5,
            'sixth': 6, 'seventh': 7, 'eighth': 8, 'ninth': 9, 'tenth': 10
        }
        
        for word, num in ordinals.items():
            if word in prompt.lower():
                return num
        
        return None

    def parse_presentation_structure(self, ai_response):
        """Parse AI response to extract presentation structure with proper formatting"""
        print(f"DEBUG: AI Response received: {ai_response[:200]}...")  # Debug info
        
        lines = ai_response.split('\n')
        slides = []
        current_slide = None
        
        for line in lines:
            line = line.strip()
            
            # Check for slide titles (starting with # or ##)
            if line.startswith('# ') or line.startswith('## '):
                # Save previous slide if it exists
                if current_slide and current_slide.get('content'):
                    slides.append(current_slide)
                    print(f"DEBUG: Added slide: {current_slide['title']} with {len(current_slide['content'])} points")
                
                # Start new slide
                title = line.replace('# ', '').replace('## ', '').strip()
                current_slide = {
                    'title': title,
                    'content': []
                }
                print(f"DEBUG: Starting new slide: {title}")
            
            # Check for bullet points
            elif line.startswith(('- ', '* ', '• ')):
                if current_slide:
                    content_point = line.replace('- ', '').replace('* ', '').replace('• ', '').strip()
                    if content_point and len(content_point) > 3:  # Only meaningful content
                        current_slide['content'].append(content_point)
                        print(f"DEBUG: Added bullet point: {content_point[:50]}...")
            
            # Handle regular text lines as content (but avoid very short lines)
            elif line and current_slide and not line.startswith('#') and len(line) > 10:
                current_slide['content'].append(line)
                print(f"DEBUG: Added regular line: {line[:50]}...")
        
        # Add the last slide
        if current_slide and current_slide.get('content'):
            slides.append(current_slide)
            print(f"DEBUG: Added final slide: {current_slide['title']} with {len(current_slide['content'])} points")
        
        print(f"DEBUG: Total slides created: {len(slides)}")
        
        # If no slides were parsed, create a default structure
        if not slides:
            print("DEBUG: No slides parsed, creating default structure")
            slides = [
                {
                    'title': 'Main Topic',
                    'content': [
                        'Key point about the topic',
                        'Important benefits and applications',
                        'Current trends and developments',
                        'Future opportunities and challenges'
                    ]
                }
            ]
        
        return slides
    
    def save_presentation(self, filename):
        """Save the current presentation"""
        if self.current_ppt:
            buffer = BytesIO()
            self.current_ppt.save(buffer)
            buffer.seek(0)
            return buffer
        return None

    def load_presentation(self, file_path_or_file):
        """Load a presentation from file path or file object"""
        try:
            from pptx import Presentation
            if isinstance(file_path_or_file, str):
                # File path
                self.current_ppt = Presentation(file_path_or_file)
            else:
                # File object
                self.current_ppt = Presentation(file_path_or_file)
            return True
        except Exception as e:
            print(f"Error loading presentation: {e}")
            return False

    def get_slide_content(self, slide_number):
        """Get content from a specific slide"""
        try:
            if not self.current_ppt or slide_number < 1:
                return None
            
            slide_index = slide_number - 1  # Convert to 0-based index
            if slide_index >= len(self.current_ppt.slides):
                return None
            
            slide = self.current_ppt.slides[slide_index]
            
            # Get title
            title = slide.shapes.title.text if slide.shapes.title else f"Slide {slide_number}"
            
            # Get content
            content = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            content.append(paragraph.text.strip())
            
            return {
                'title': title,
                'content': content
            }
        except Exception as e:
            print(f"Error getting slide content: {e}")
            return None

    def get_presentation_summary(self):
        """Get a summary of the current presentation"""
        if not self.current_ppt:
            return None
        
        try:
            slides_info = []
            for i, slide in enumerate(self.current_ppt.slides, 1):
                slide_info = {
                    'number': i,
                    'title': 'Untitled Slide',
                    'content_points': 0
                }
                
                # Get title
                if slide.shapes.title and slide.shapes.title.text:
                    slide_info['title'] = slide.shapes.title.text.strip()
                
                # Count content points
                content_count = 0
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                content_count += 1
                
                slide_info['content_points'] = content_count
                slides_info.append(slide_info)
            
            return {
                'total_slides': len(self.current_ppt.slides),
                'slides': slides_info
            }
        except Exception as e:
            print(f"Error getting presentation summary: {e}")
            return None

def main():
    st.set_page_config(page_title="PowerPoint AI Chatbot", layout="wide")
    
    st.title("🤖 PowerPoint AI Chatbot")
    st.markdown("Create, edit, and enhance PowerPoint presentations with AI assistance!")
    
    # Initialize chatbot
    if 'chatbot' not in st.session_state:
        st.session_state.chatbot = PowerPointChatbot()
    
    # Sidebar for options
    st.sidebar.title("Options")
    
    # Initialize operation in session state if not exists
    if 'operation' not in st.session_state:
        st.session_state['operation'] = "Create New Presentation"
    
    operation = st.sidebar.selectbox(
        "Choose operation:",
        ["Create New Presentation", "Upload & Edit PPT", "Add Chart/Visualization", "Chat with AI"],
        index=["Create New Presentation", "Upload & Edit PPT", "Add Chart/Visualization", "Chat with AI"].index(st.session_state['operation'])
    )
    
    # Update session state when selection changes
    if operation != st.session_state['operation']:
        st.session_state['operation'] = operation
    
    if operation == "Create New Presentation":
        st.header("Create New Presentation")
        
        col1, col2 = st.columns([2, 1])
        
        with col1:
            presentation_topic = st.text_input("Presentation Topic:", 
                                             placeholder="e.g., Digital Marketing Strategy")
            
            additional_requirements = st.text_area(
                "Additional Requirements:",
                placeholder="e.g., Include 5 slides, focus on social media, add statistics"
            )
            
            if st.button("Generate Presentation"):
                if presentation_topic:
                    with st.spinner("Generating presentation structure..."):
                        prompt = f"""
                        Create a detailed and professional presentation structure about "{presentation_topic}".
                        {additional_requirements if additional_requirements else ""}
                        
                        Requirements:
                        - Create 4-6 slides with clear, concise content
                        - Each slide should have 3-5 bullet points maximum
                        - Each bullet point should be ONE clear, complete sentence (not paragraphs)
                        - Use professional language and terminology
                        - Make each bullet point specific and actionable
                        - Keep bullet points between 10-25 words each
                        
                        Format your response EXACTLY as follows (use this exact format):
                        
                        # Introduction to {presentation_topic}
                        - Brief overview of the topic and its significance
                        - Key objectives and goals of this presentation
                        - Why this topic is relevant and important today
                        - Current market landscape and opportunities
                        
                        # Key Concepts and Fundamentals
                        - Core principles and foundational concepts
                        - Important terminology and definitions
                        - Historical context and background information
                        - Current state of the industry or field
                        
                        # Applications and Implementation
                        - Real-world applications and use cases
                        - Practical implementation strategies
                        - Success stories and case studies
                        - Benefits and advantages for organizations
                        
                        # Challenges and Solutions
                        - Common obstacles and potential problems
                        - Proven solutions and best practices
                        - Risk mitigation strategies
                        - Lessons learned from industry experts
                        
                        # Future Outlook and Trends
                        - Emerging trends and developments
                        - Future opportunities and growth potential
                        - Predicted industry changes and evolution
                        - Strategic recommendations for stakeholders
                        
                        # Conclusion and Next Steps
                        - Summary of key takeaways and insights
                        - Actionable recommendations for audience
                        - Next steps for implementation or exploration
                        - Call to action and follow-up resources
                        
                        IMPORTANT: Each bullet point must be a single, clear sentence. Do not write paragraphs or multiple sentences in one bullet point.
                        """
                        
                        ai_response = st.session_state.chatbot.generate_content_with_ai(prompt)
                        st.write("AI Generated Structure:")
                        st.write(ai_response)
                        
                        # Parse and create presentation
                        slides_structure = st.session_state.chatbot.parse_presentation_structure(ai_response)
                        
                        if slides_structure:
                            prs = st.session_state.chatbot.create_presentation(
                                presentation_topic, slides_structure
                            )
                            
                            st.success("Presentation created successfully!")
                            
                            # Download button
                            buffer = st.session_state.chatbot.save_presentation("presentation.pptx")
                            if buffer:
                                st.download_button(
                                    label="📥 Download Presentation",
                                    data=buffer.getvalue(),
                                    file_name=f"{presentation_topic.replace(' ', '_')}.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                )
        
        with col2:
            st.info("💡 **Tips:**\n\n- Be specific about your topic\n- Mention the number of slides you want\n- Include any specific requirements\n- You can edit slides after creation")
    
    elif operation == "Upload & Edit PPT":
        st.header("Upload & Edit PowerPoint")
        
        uploaded_file = st.file_uploader("Upload PowerPoint file", type=['pptx'])
        
        if uploaded_file is not None:
            # Load the uploaded presentation
            prs = Presentation(uploaded_file)
            st.session_state.chatbot.current_ppt = prs
            
            st.success("Presentation loaded successfully!")
            
            # Display slides for editing
            st.subheader("Edit Slides")
            
            slide_count = len(prs.slides)
            slide_to_edit = st.selectbox("Select slide to edit:", range(slide_count))
            
            # Show current slide content
            try:
                slide = prs.slides[slide_to_edit]
                st.write(f"**Current Slide {slide_to_edit + 1} Title:** {slide.shapes.title.text}")
                
                # Extract current content
                current_content = []
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                current_content.append(paragraph.text)
                
                st.write("**Current Content:**")
                for content in current_content:
                    st.write(f"- {content}")
                
                # Edit form
                with st.form("edit_slide_form"):
                    new_title = st.text_input("New Title:", value=slide.shapes.title.text)
                    new_content_text = st.text_area(
                        "New Content (one point per line):",
                        value="\n".join(current_content)
                    )
                    
                    submitted = st.form_submit_button("Update Slide")
                
                # Handle form submission outside the form
                if submitted:
                    new_content = {
                        'title': new_title,
                        'content': new_content_text.split('\n') if new_content_text else []
                    }
                    
                    success = st.session_state.chatbot.edit_slide_content(slide_to_edit, new_content)
                    
                    if success:
                        st.success("Slide updated successfully!")
                        
                        # Download updated presentation (moved outside form)
                        buffer = st.session_state.chatbot.save_presentation("updated_presentation.pptx")
                        if buffer:
                            st.download_button(
                                label="📥 Download Updated Presentation",
                                data=buffer.getvalue(),
                                file_name="updated_presentation.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                    else:
                        st.error("Failed to update slide.")
            
            except Exception as e:
                st.error(f"Error reading slide: {str(e)}")
    
    elif operation == "Add Chart/Visualization":
        st.header("Add Chart/Visualization")
        
        if st.session_state.chatbot.current_ppt is None:
            st.warning("Please create or upload a presentation first.")
        else:
            chart_type = st.selectbox("Chart Type:", ["bar", "line", "pie"])
            
            # Data input methods
            data_input_method = st.radio("Data Input Method:", ["Manual Entry", "Upload CSV"])
            
            if data_input_method == "Manual Entry":
                st.subheader("Enter Chart Data")
                
                chart_title = st.text_input("Chart Title:", "Sales Data")
                
                # Simple data entry
                labels_input = st.text_input("Labels (comma-separated):", "Q1,Q2,Q3,Q4")
                values_input = st.text_input("Values (comma-separated):", "100,150,120,180")
                
                if st.button("Add Chart to Presentation"):
                    try:
                        labels = [label.strip() for label in labels_input.split(',')]
                        values = [float(value.strip()) for value in values_input.split(',')]
                        
                        chart_data = {
                            'title': chart_title,
                            'labels': labels,
                            'values': values
                        }
                        
                        success = st.session_state.chatbot.add_chart_slide(chart_data, chart_type)
                        
                        if success:
                            st.success("Chart added to presentation!")
                            
                            # Download updated presentation
                            buffer = st.session_state.chatbot.save_presentation("presentation_with_chart.pptx")
                            if buffer:
                                st.download_button(
                                    label="📥 Download Presentation with Chart",
                                    data=buffer.getvalue(),
                                    file_name="presentation_with_chart.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                )
                        else:
                            st.error("Failed to add chart.")
                    
                    except ValueError:
                        st.error("Please enter valid numeric values.")
            
            elif data_input_method == "Upload CSV":
                uploaded_csv = st.file_uploader("Upload CSV file", type=['csv'])
                
                if uploaded_csv is not None:
                    df = pd.read_csv(uploaded_csv)
                    st.write("Data Preview:")
                    st.dataframe(df.head())
                    
                    # Column selection
                    label_column = st.selectbox("Select label column:", df.columns)
                    value_column = st.selectbox("Select value column:", df.columns)
                    chart_title = st.text_input("Chart Title:", "Data Visualization")
                    
                    if st.button("Create Chart from CSV"):
                        chart_data = {
                            'title': chart_title,
                            'labels': df[label_column].tolist(),
                            'values': df[value_column].tolist()
                        }
                        
                        success = st.session_state.chatbot.add_chart_slide(chart_data, chart_type)
                        
                        if success:
                            st.success("Chart created from CSV and added to presentation!")
                            
                            buffer = st.session_state.chatbot.save_presentation("presentation_with_csv_chart.pptx")
                            if buffer:
                                st.download_button(
                                    label="📥 Download Presentation",
                                    data=buffer.getvalue(),
                                    file_name="presentation_with_csv_chart.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                )
    
    elif operation == "Chat with AI":
        st.header("Chat with AI Assistant")
        
        if "messages" not in st.session_state:
            st.session_state.messages = []
        
        # File upload section state
        if 'show_upload' not in st.session_state:
            st.session_state.show_upload = False
        
        # Create containers for proper layout
        message_container = st.container()
        upload_container = st.container()
        input_container = st.container()
        
        # Handle the input first (but display at bottom)
        with input_container:
            # Add some spacing before the input section
            st.markdown("<br>", unsafe_allow_html=True)
            
            # Create the bottom input area using columns
            col_input, col_button = st.columns([9, 1])
            
            with col_input:
                prompt = st.chat_input("💬 Ask me anything about your presentation...")
            
            with col_button:
                if st.button("➕", help="Upload PowerPoint file", key="upload_toggle", use_container_width=True):
                    st.session_state.show_upload = not st.session_state.show_upload
                    st.rerun()
        
        # Process prompt immediately if entered
        if prompt:
            st.session_state.messages.append({"role": "user", "content": prompt})
            
            # Check if user is asking for presentation creation
            presentation_keywords = [
                "create a presentation", "make a presentation", "generate a presentation",
                "create presentation", "make presentation", "generate presentation",
                "create ppt", "make ppt", "generate ppt", "create powerpoint",
                "presentation on", "presentation about", "slides about", "slides on"
            ]
            
            is_presentation_request = any(keyword in prompt.lower() for keyword in presentation_keywords)
            
            # Check if user is asking for slide editing
            editing_keywords = [
                "edit slide", "modify slide", "change slide", "update slide",
                "edit content", "modify content", "change content", "update content",
                "edit presentation", "modify presentation", "update presentation"
            ]
            
            # Check if user is asking to add new slides
            add_slide_keywords = [
                "add slide", "add new slide", "create slide", "insert slide",
                "add a slide", "create new slide", "new slide about"
            ]
            
            # Check if user is asking to view slide content
            view_slide_keywords = [
                "show me slide", "view slide", "display slide", "what's in slide",
                "slide content", "content of slide", "show slide"
            ]
            
            is_editing_request = any(keyword in prompt.lower() for keyword in editing_keywords)
            is_add_slide_request = any(keyword in prompt.lower() for keyword in add_slide_keywords)
            is_view_slide_request = any(keyword in prompt.lower() for keyword in view_slide_keywords)
            
            # Generate AI response and add to messages
            if is_presentation_request:
                # Extract topic from the prompt
                topic = st.session_state.chatbot.extract_topic_from_prompt(prompt)
                
                # Extract number of slides from user request
                slide_count = st.session_state.chatbot.extract_slide_count_from_prompt(prompt)
                
                # Generate presentation structure using AI
                structure_prompt = f"""
                Create a detailed and professional presentation structure about "{topic}".
                Extract any specific requirements from this user request: "{prompt}"
                
                Requirements:
                - Create EXACTLY {slide_count} slides with clear, concise content
                - Each slide should have 3-5 bullet points maximum
                - Each bullet point should be ONE clear, complete sentence (not paragraphs)
                - Use professional language and terminology
                - Make each bullet point specific and actionable
                - Keep bullet points between 10-25 words each
                
                Generate EXACTLY {slide_count} slides following this structure pattern:

                # Slide Title 1
                - Bullet point 1 (10-25 words)
                - Bullet point 2 (10-25 words)
                - Bullet point 3 (10-25 words)
                - Bullet point 4 (10-25 words)

                # Slide Title 2
                - Bullet point 1 (10-25 words)
                - Bullet point 2 (10-25 words)
                - Bullet point 3 (10-25 words)
                - Bullet point 4 (10-25 words)

                Continue this pattern for all {slide_count} slides. Make sure each slide has a clear, descriptive title and 3-5 bullet points.
                """
                
                try:
                    response = st.session_state.chatbot.model.generate_content(structure_prompt)
                    
                    # Parse the response and create the presentation
                    content_structure = st.session_state.chatbot.parse_presentation_structure(response.text)
                    
                    if content_structure:
                        # Create the presentation
                        prs = st.session_state.chatbot.create_presentation(topic, content_structure)
                        
                        # Store the presentation in the chatbot for persistence
                        st.session_state.chatbot.current_ppt = prs
                        
                        # Convert to bytes for download
                        pptx_buffer = BytesIO()
                        prs.save(pptx_buffer)
                        pptx_data = pptx_buffer.getvalue()
                        
                        # Create response message
                        response_text = f"🎯 I've created a **{slide_count}-slide presentation** about **{topic}**!\n\n"
                        response_text += f"**📊 Presentation Overview:**\n"
                        for i, slide_data in enumerate(content_structure, 1):
                            response_text += f"• **Slide {i}:** {slide_data['title']}\n"
                        
                        response_text += f"\n**✅ Your presentation is ready!** Click the download button below to get your PowerPoint file."
                        
                        # Add AI response to messages
                        st.session_state.messages.append({"role": "assistant", "content": response_text})
                        
                        # Store the presentation data for download with error checking
                        try:
                            st.session_state.pptx_data = pptx_data
                            # Sanitize filename by removing unsafe characters
                            safe_topic = topic.replace(' ', '_').replace(':', '').replace('\n', '').replace('\r', '').replace('/', '_').replace('\\', '_').replace('?', '').replace('*', '').replace('<', '').replace('>', '').replace('|', '').replace('"', '').lower()
                            st.session_state.filename = f"{safe_topic}_presentation.pptx"
                        except Exception as storage_error:
                            st.error(f"Error storing presentation data: {str(storage_error)}")
                        
                    else:
                        error_msg = "❌ I couldn't create the presentation. Could you please rephrase your request?"
                        st.session_state.messages.append({"role": "assistant", "content": error_msg})
                        
                except Exception as e:
                    error_msg = f"❌ Error creating presentation: {str(e)}"
                    st.session_state.messages.append({"role": "assistant", "content": error_msg})
            
            elif is_editing_request:
                if st.session_state.chatbot.current_ppt is None:
                    response_text = "📎 **Please upload a PowerPoint file first** before editing.\n\nClick the ➕ button to upload your presentation!"
                    st.session_state.messages.append({"role": "assistant", "content": response_text})
                else:
                    # Process editing request
                    try:
                        # Extract slide number from prompt
                        slide_number = st.session_state.chatbot.extract_slide_number_from_prompt(prompt)
                        
                        if slide_number:
                            # Get current slide content
                            current_slide = st.session_state.chatbot.get_slide_content(slide_number)
                            
                            if current_slide:
                                # Generate AI prompt for editing
                                edit_prompt = f"""
                                The user wants to edit slide {slide_number} of their presentation.
                                
                                Current slide content:
                                Title: {current_slide['title']}
                                Content: {', '.join(current_slide['content'])}
                                
                                User's editing request: "{prompt}"
                                
                                Generate new content for this slide based on the user's request. 
                                Format your response as:
                                
                                # New Slide Title
                                - Bullet point 1
                                - Bullet point 2
                                - Bullet point 3
                                - Bullet point 4
                                
                                Keep bullet points concise (10-25 words each) and professional.
                                """
                                
                                ai_response = st.session_state.chatbot.model.generate_content(edit_prompt)
                                
                                # Parse the AI response to get new slide content
                                new_slide_structure = st.session_state.chatbot.parse_presentation_structure(ai_response.text)
                                
                                if new_slide_structure and len(new_slide_structure) > 0:
                                    new_slide_content = new_slide_structure[0]
                                    
                                    # Update the slide
                                    success = st.session_state.chatbot.edit_slide_content(slide_number - 1, new_slide_content)
                                    
                                    if success:
                                        # Create updated presentation for download
                                        pptx_buffer = BytesIO()
                                        st.session_state.chatbot.current_ppt.save(pptx_buffer)
                                        pptx_data = pptx_buffer.getvalue()
                                        
                                        # Store updated presentation
                                        st.session_state.pptx_data = pptx_data
                                        st.session_state.filename = "updated_presentation.pptx"
                                        
                                        response_text = f"✅ **Successfully updated slide {slide_number}!**\n\n"
                                        response_text += f"**🔄 Changes made:**\n"
                                        response_text += f"• **New Title:** {new_slide_content['title']}\n"
                                        response_text += f"• **Updated Content:** {len(new_slide_content['content'])} bullet points\n\n"
                                        response_text += f"**📥 Your updated presentation is ready!** Click the download button below."
                                        
                                        st.session_state.messages.append({"role": "assistant", "content": response_text})
                                    else:
                                        st.session_state.messages.append({"role": "assistant", "content": "❌ Failed to update the slide. Please try again."})
                                else:
                                    st.session_state.messages.append({"role": "assistant", "content": "❌ Couldn't parse the editing request. Please be more specific about what you want to change."})
                            else:
                                st.session_state.messages.append({"role": "assistant", "content": f"❌ Slide {slide_number} not found. Please check the slide number."})
                        else:
                            # General editing request without specific slide number
                            slide_count = len(st.session_state.chatbot.current_ppt.slides)
                            response_text = f"📝 **I can help you edit your presentation!**\n\n"
                            response_text += f"**📊 Your presentation has {slide_count} slides.**\n\n"
                            response_text += f"**💡 Try commands like:**\n"
                            response_text += f"• *'Edit slide 2 title to New Marketing Strategy'*\n"
                            response_text += f"• *'Modify slide 3 content about social media'*\n"
                            response_text += f"• *'Change slide 1 to focus on digital transformation'*\n\n"
                            response_text += f"**Which slide would you like to edit?**"
                            
                            st.session_state.messages.append({"role": "assistant", "content": response_text})
                            
                    except Exception as e:
                        error_msg = f"❌ Error processing edit request: {str(e)}"
                        st.session_state.messages.append({"role": "assistant", "content": error_msg})
            
            elif is_add_slide_request:
                if st.session_state.chatbot.current_ppt is None:
                    response_text = "📎 **Please upload a PowerPoint file first** before adding slides.\n\nClick the ➕ button to upload your presentation!"
                    st.session_state.messages.append({"role": "assistant", "content": response_text})
                else:
                    # Process add slide request
                    try:
                        # Generate new slide content using AI
                        slide_prompt = f"""
                        The user wants to add a new slide to their presentation.
                        User's request: "{prompt}"
                        
                        Create content for this new slide based on the user's request.
                        
                        Format your response as:
                        
                        # Slide Title (extract from user's request)
                        - Bullet point 1 (10-25 words)
                        - Bullet point 2 (10-25 words)
                        - Bullet point 3 (10-25 words)
                        - Bullet point 4 (10-25 words)
                        
                        Keep bullet points concise, professional, and relevant to the topic.
                        """
                        
                        ai_response = st.session_state.chatbot.model.generate_content(slide_prompt)
                        
                        # Parse the AI response to get new slide content
                        new_slide_structure = st.session_state.chatbot.parse_presentation_structure(ai_response.text)
                        
                        if new_slide_structure and len(new_slide_structure) > 0:
                            new_slide_content = new_slide_structure[0]
                            
                            # Add the new slide
                            success = st.session_state.chatbot.add_new_slide(new_slide_content)
                            
                            if success:
                                # Get current slide count
                                slide_count = len(st.session_state.chatbot.current_ppt.slides)
                                
                                # Create updated presentation for download
                                pptx_buffer = BytesIO()
                                st.session_state.chatbot.current_ppt.save(pptx_buffer)
                                pptx_data = pptx_buffer.getvalue()
                                
                                # Store updated presentation
                                st.session_state.pptx_data = pptx_data
                                st.session_state.filename = "updated_presentation.pptx"
                                
                                response_text = f"✅ **Successfully added new slide!**\n\n"
                                response_text += f"**📊 New Slide Details:**\n"
                                response_text += f"• **Position:** Slide {slide_count} (last slide)\n"
                                response_text += f"• **Title:** {new_slide_content['title']}\n"
                                response_text += f"• **Content:** {len(new_slide_content['content'])} bullet points\n\n"
                                response_text += f"**📈 Your presentation now has {slide_count} slides total.**\n\n"
                                response_text += f"**📥 Download your updated presentation below!**"
                                
                                st.session_state.messages.append({"role": "assistant", "content": response_text})
                            else:
                                st.session_state.messages.append({"role": "assistant", "content": "❌ Failed to add the new slide. Please try again."})
                        else:
                            st.session_state.messages.append({"role": "assistant", "content": "❌ Couldn't understand what content to add to the new slide. Please be more specific."})
                            
                    except Exception as e:
                        error_msg = f"❌ Error adding new slide: {str(e)}"
                        st.session_state.messages.append({"role": "assistant", "content": error_msg})
            
            elif is_view_slide_request:
                if st.session_state.chatbot.current_ppt is None:
                    response_text = "📎 **Please upload a PowerPoint file first** before viewing slides.\n\nClick the ➕ button to upload your presentation!"
                    st.session_state.messages.append({"role": "assistant", "content": response_text})
                else:
                    try:
                        # Extract slide number from prompt
                        slide_number = st.session_state.chatbot.extract_slide_number_from_prompt(prompt)
                        
                        if slide_number:
                            # Get slide content
                            slide_content = st.session_state.chatbot.get_slide_content(slide_number)
                            
                            if slide_content:
                                response_text = f"📖 **Slide {slide_number} Content:**\n\n"
                                response_text += f"**🏷️ Title:** {slide_content['title']}\n\n"
                                
                                if slide_content['content']:
                                    response_text += f"**📝 Content:**\n"
                                    for i, point in enumerate(slide_content['content'], 1):
                                        clean_point = point.replace('•', '').replace('-', '').replace('*', '').strip()
                                        response_text += f"{i}. {clean_point}\n"
                                else:
                                    response_text += f"**📝 Content:** No content found\n"
                                
                                response_text += f"\n**💡 Want to edit this slide?** Try:\n"
                                response_text += f"*'Edit slide {slide_number} title to [new title]'*\n"
                                response_text += f"*'Modify slide {slide_number} content about [topic]'*"
                                
                                st.session_state.messages.append({"role": "assistant", "content": response_text})
                            else:
                                st.session_state.messages.append({"role": "assistant", "content": f"❌ Slide {slide_number} not found. Please check the slide number."})
                        else:
                            # Show overview of all slides
                            presentation_summary = st.session_state.chatbot.get_presentation_summary()
                            
                            if presentation_summary:
                                response_text = f"📊 **Presentation Overview:**\n\n"
                                response_text += f"**Total Slides:** {presentation_summary['total_slides']}\n\n"
                                response_text += f"**📋 All Slides:**\n"
                                
                                for slide_info in presentation_summary['slides']:
                                    title = slide_info['title'][:40] + "..." if len(slide_info['title']) > 40 else slide_info['title']
                                    response_text += f"• **Slide {slide_info['number']}:** {title} ({slide_info['content_points']} points)\n"
                                
                                response_text += f"\n**💡 To view specific slide:** *'Show me slide [number]'*"
                                st.session_state.messages.append({"role": "assistant", "content": response_text})
                            else:
                                st.session_state.messages.append({"role": "assistant", "content": "❌ Could not read presentation structure."})
                                
                    except Exception as e:
                        error_msg = f"❌ Error viewing slide: {str(e)}"
                        st.session_state.messages.append({"role": "assistant", "content": error_msg})
            
            else:
                # General chat response
                try:
                    chat_prompt = f"""
                    You are a PowerPoint presentation assistant. The user said: "{prompt}"
                    
                    Provide a helpful response about PowerPoint presentations, slide creation, or presentation tips.
                    If they're asking about creating presentations, guide them to use phrases like "create a presentation about [topic]".
                    If they're asking about editing, guide them to upload a file first using the + button.
                    
                    Keep your response friendly, professional, and focused on PowerPoint assistance.
                    """
                    
                    response = st.session_state.chatbot.model.generate_content(chat_prompt)
                    st.session_state.messages.append({"role": "assistant", "content": response.text})
                    
                except Exception as e:
                    error_msg = f"❌ Error generating response: {str(e)}"
                    st.session_state.messages.append({"role": "assistant", "content": error_msg})
            
            # Rerun to update the display
            st.rerun()
        
        # Display all messages in the message container
        with message_container:
            for message in st.session_state.messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])
                    
                    # Show download button for presentations
                    if message["role"] == "assistant" and any(phrase in message["content"] for phrase in [
                        "Your presentation is ready!", 
                        "Your updated presentation is ready!", 
                        "Download your updated presentation below!",
                        "📥 Download your updated presentation below!"
                    ]):
                        if hasattr(st.session_state, 'pptx_data') and st.session_state.pptx_data:
                            try:
                                import time
                                download_key = f"chat_download_{int(time.time())}_{hash(message['content']) % 10000}"
                                
                                st.download_button(
                                    label="📥 Download PowerPoint Presentation",
                                    data=st.session_state.pptx_data,
                                    file_name=getattr(st.session_state, 'filename', 'presentation.pptx'),
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    key=download_key,
                                    use_container_width=True
                                )
                            except Exception as e:
                                st.error(f"Download error: {str(e)}")
                                # Fallback: Try to recreate the presentation data
                                if hasattr(st.session_state.chatbot, 'current_ppt') and st.session_state.chatbot.current_ppt:
                                    try:
                                        fallback_buffer = BytesIO()
                                        st.session_state.chatbot.current_ppt.save(fallback_buffer)
                                        fallback_data = fallback_buffer.getvalue()
                                        
                                        st.download_button(
                                            label="📥 Download Presentation (Fallback)",
                                            data=fallback_data,
                                            file_name="updated_presentation.pptx",
                                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                            key=f"fallback_{download_key}",
                                            use_container_width=True
                                        )
                                    except Exception as fe:
                                        st.error(f"Fallback download also failed: {str(fe)}")
        
        # Show upload section in the upload container
        with upload_container:
            if st.session_state.show_upload:
                st.markdown("### 📎 Upload PowerPoint File")
                uploaded_file = st.file_uploader(
                    "Choose a PowerPoint file to edit:",
                    type=['pptx'],
                    help="Upload your PowerPoint presentation to edit it",
                    key="main_file_upload"
                )
                
                if uploaded_file is not None:
                    with st.spinner(f"Loading {uploaded_file.name}..."):
                        try:
                            # Save uploaded file temporarily
                            temp_path = f"temp_{uploaded_file.name}"
                            with open(temp_path, "wb") as f:
                                f.write(uploaded_file.getbuffer())
                            
                            # Load the presentation
                            success = st.session_state.chatbot.load_presentation(temp_path)
                            
                            if success:
                                # Add file upload message to chat
                                upload_message = f"📎 Uploaded presentation: {uploaded_file.name}"
                                st.session_state.messages.append({"role": "user", "content": upload_message})
                                
                                # Get presentation summary
                                presentation_summary = st.session_state.chatbot.get_presentation_summary()
                                
                                # Add AI response about the uploaded file
                                ai_response = f"✅ **Successfully loaded '{uploaded_file.name}'!**\n\n"
                                
                                if presentation_summary:
                                    ai_response += f"**📊 Presentation Overview:**\n"
                                    ai_response += f"• **Total Slides:** {presentation_summary['total_slides']}\n\n"
                                    ai_response += f"**📋 Slide Summary:**\n"
                                    
                                    for slide_info in presentation_summary['slides'][:5]:  # Show first 5 slides
                                        title = slide_info['title'][:50] + "..." if len(slide_info['title']) > 50 else slide_info['title']
                                        ai_response += f"• **Slide {slide_info['number']}:** {title} ({slide_info['content_points']} points)\n"
                                    
                                    if presentation_summary['total_slides'] > 5:
                                        ai_response += f"• *...and {presentation_summary['total_slides'] - 5} more slides*\n"
                                
                                ai_response += f"\n**💡 What you can do now:**\n"
                                ai_response += f"• **Edit content:** *'Edit slide 2 title to New Marketing Strategy'*\n"
                                ai_response += f"• **Add slides:** *'Add a new slide about market analysis'*\n"
                                ai_response += f"• **View content:** *'Show me slide 3'*\n"
                                ai_response += f"• **Modify multiple slides:** *'Update slide 1 and 2 with latest data'*\n\n"
                                ai_response += f"**🚀 Ready to help! What would you like to do?**"
                                
                                st.session_state.messages.append({"role": "assistant", "content": ai_response})
                                
                                # Clean up temp file
                                import os
                                if os.path.exists(temp_path):
                                    os.remove(temp_path)
                                
                                # Hide upload section after successful upload
                                st.session_state.show_upload = False
                                st.success("File uploaded successfully!")
                                st.rerun()
                            else:
                                st.error("❌ Failed to load the presentation. Please ensure it's a valid PowerPoint file.")
                        
                        except Exception as e:
                            st.error(f"❌ Error uploading file: {str(e)}")
                            # Clean up temp file if it exists
                            import os
                            if os.path.exists(temp_path):
                                os.remove(temp_path)

if __name__ == "__main__":
    main()